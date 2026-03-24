"""Office file conversion: PPTX → PNG slides, DOCX → HTML."""

import base64
import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# ── PPTX to Images ───────────────────────────────────────────────

def pptx_to_images(pptx_path: str, output_dir: Path) -> list[Path]:
    """
    Extract slides from a PPTX and render them as PNG images.

    This is a best-effort renderer using python-pptx and Pillow.
    It handles text frames, embedded pictures, and tables.
    Complex elements (charts, SmartArt, grouped shapes) are rendered
    as placeholder rectangles.

    Returns a list of paths to the generated PNG files.
    """
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(pptx_path)

    # Convert slide dimensions from EMU to pixels at 96 DPI
    dpi = 150  # Higher DPI for better quality
    emu_per_inch = 914400
    slide_w = int(prs.slide_width / emu_per_inch * dpi)
    slide_h = int(prs.slide_height / emu_per_inch * dpi)

    output_paths = []

    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (slide_w, slide_h), "white")
        draw = ImageDraw.Draw(img)

        # Try to render slide background
        _render_slide_background(slide, img, slide_w, slide_h)

        for shape in slide.shapes:
            try:
                if hasattr(shape, "image") and shape.shape_type is not None:
                    _render_picture(img, shape, slide_w, slide_h, prs, dpi)
                elif shape.has_text_frame:
                    _render_text_frame(draw, shape, slide_w, slide_h, prs, dpi)
                elif shape.has_table:
                    _render_table(draw, shape, slide_w, slide_h, prs, dpi)
                else:
                    _render_placeholder(draw, shape, slide_w, slide_h, prs, dpi)
            except Exception as e:
                logger.debug(f"Slide {idx}: skipped shape '{shape.name}': {e}")

        out_path = output_dir / f"slide_{idx:04d}.png"
        img.save(str(out_path), "PNG")
        output_paths.append(out_path)
        logger.debug(f"Rendered slide {idx} -> {out_path}")

    return output_paths


def _emu_to_px(emu_val, dpi=150):
    """Convert EMU value to pixels."""
    return int(emu_val / 914400 * dpi)


def _get_font(size_pt=12, bold=False):
    """Try to load a suitable font, fall back to default."""
    from PIL import ImageFont
    size = max(8, int(size_pt))
    try:
        if bold:
            return ImageFont.truetype("arialbd.ttf", size)
        return ImageFont.truetype("arial.ttf", size)
    except (OSError, IOError):
        try:
            return ImageFont.truetype("segoeui.ttf", size)
        except (OSError, IOError):
            return ImageFont.load_default()


def _render_slide_background(slide, img, w, h):
    """Attempt to render the slide background color."""
    from PIL import ImageDraw
    try:
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            color = bg.fill.fore_color.rgb
            hex_color = f"#{color}"
            draw = ImageDraw.Draw(img)
            draw.rectangle([0, 0, w, h], fill=hex_color)
    except Exception:
        pass  # Keep white background


def _render_text_frame(draw, shape, slide_w, slide_h, prs, dpi):
    """Render text content from a shape's text frame."""
    x = _emu_to_px(shape.left, dpi)
    y = _emu_to_px(shape.top, dpi)

    for para in shape.text_frame.paragraphs:
        line_text = ""
        font_size = 12
        is_bold = False

        for run in para.runs:
            line_text += run.text
            if run.font.size:
                font_size = run.font.size.pt * (dpi / 96)
            if run.font.bold:
                is_bold = True

        if line_text.strip():
            font = _get_font(font_size, is_bold)

            # Determine text color
            fill = "black"
            try:
                if para.runs and para.runs[0].font.color.rgb:
                    fill = f"#{para.runs[0].font.color.rgb}"
            except Exception:
                pass

            draw.text((x, y), line_text, fill=fill, font=font)

        y += int(font_size * 1.4) + 2


def _render_picture(img, shape, slide_w, slide_h, prs, dpi):
    """Render an embedded picture shape."""
    from PIL import Image as PILImage

    x = _emu_to_px(shape.left, dpi)
    y = _emu_to_px(shape.top, dpi)
    w = _emu_to_px(shape.width, dpi)
    h = _emu_to_px(shape.height, dpi)

    blob = shape.image.blob
    pic = PILImage.open(io.BytesIO(blob))
    pic = pic.convert("RGBA")
    pic = pic.resize((max(1, w), max(1, h)), PILImage.LANCZOS)
    img.paste(pic, (x, y), pic if pic.mode == "RGBA" else None)


def _render_table(draw, shape, slide_w, slide_h, prs, dpi):
    """Render a table shape with cell text."""
    x0 = _emu_to_px(shape.left, dpi)
    y0 = _emu_to_px(shape.top, dpi)
    table = shape.table

    row_count = len(table.rows)
    col_count = len(table.columns)
    if row_count == 0 or col_count == 0:
        return

    total_w = _emu_to_px(shape.width, dpi)
    total_h = _emu_to_px(shape.height, dpi)
    cell_w = total_w // col_count
    cell_h = total_h // row_count

    font = _get_font(10)

    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cx = x0 + c_idx * cell_w
            cy = y0 + r_idx * cell_h
            draw.rectangle([cx, cy, cx + cell_w, cy + cell_h], outline="#999999")
            text = cell.text.strip()
            if text:
                draw.text((cx + 4, cy + 4), text, fill="black", font=font)


def _render_placeholder(draw, shape, slide_w, slide_h, prs, dpi):
    """Draw a placeholder rectangle for unsupported shape types."""
    try:
        x = _emu_to_px(shape.left, dpi)
        y = _emu_to_px(shape.top, dpi)
        w = _emu_to_px(shape.width, dpi)
        h = _emu_to_px(shape.height, dpi)
        if w > 10 and h > 10:
            draw.rectangle([x, y, x + w, y + h], outline="#cccccc", fill="#f5f5f5")
    except Exception:
        pass


# ── DOCX to HTML ─────────────────────────────────────────────────

def docx_to_html(docx_path: str) -> str:
    """
    Convert a DOCX file to a self-contained HTML document.

    Handles headings, paragraphs, bold/italic/underline formatting,
    inline images (base64-encoded), and tables.
    """
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    doc = Document(docx_path)

    parts = [
        '<!DOCTYPE html>',
        '<html><head><meta charset="utf-8">',
        '<style>',
        'body { font-family: "Segoe UI", Calibri, Arial, sans-serif; ',
        '  padding: 50px 60px; max-width: 900px; margin: 0 auto; ',
        '  background: #ffffff; color: #333; line-height: 1.6; }',
        'h1 { font-size: 2em; color: #1a1a2e; margin-top: 0.8em; }',
        'h2 { font-size: 1.6em; color: #2d2d44; }',
        'h3 { font-size: 1.3em; color: #3d3d5c; }',
        'h4, h5, h6 { font-size: 1.1em; color: #555; }',
        'p { margin: 0.5em 0; }',
        'table { border-collapse: collapse; width: 100%; margin: 1em 0; }',
        'td, th { border: 1px solid #d0d0d0; padding: 8px 12px; text-align: left; }',
        'th { background: #f0f0f5; font-weight: bold; }',
        'tr:nth-child(even) { background: #fafafa; }',
        'img { max-width: 100%; height: auto; margin: 1em 0; }',
        'ul, ol { padding-left: 2em; }',
        '</style>',
        '</head><body>',
    ]

    # Build a map of relationship IDs to image blobs
    image_map = {}
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_map[rel.rId] = rel.target_part.blob
    except Exception as e:
        logger.debug(f"Could not extract images: {e}")

    # Process paragraphs
    for para in doc.paragraphs:
        tag = _style_to_tag(para.style.name)
        inner = _runs_to_html(para.runs)

        # Check for inline images in the paragraph XML
        img_html = _extract_inline_images(para, image_map)
        if img_html:
            inner += img_html

        if inner.strip():
            parts.append(f'<{tag}>{inner}</{tag}>')

    # Process tables
    for table in doc.tables:
        parts.append(_table_to_html(table))

    parts.append('</body></html>')
    return '\n'.join(parts)


def _style_to_tag(style_name: str) -> str:
    """Map Word paragraph style names to HTML tags."""
    name_lower = style_name.lower()
    for i in range(1, 7):
        if f"heading {i}" in name_lower:
            return f"h{i}"
    if "title" in name_lower:
        return "h1"
    if "subtitle" in name_lower:
        return "h2"
    if "list" in name_lower:
        return "li"
    return "p"


def _runs_to_html(runs) -> str:
    """Convert a list of docx runs to HTML with inline formatting."""
    html = ""
    for run in runs:
        text = run.text
        if not text:
            continue
        # Escape HTML entities
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if run.bold:
            text = f"<strong>{text}</strong>"
        if run.italic:
            text = f"<em>{text}</em>"
        if run.underline:
            text = f"<u>{text}</u>"
        html += text
    return html


def _extract_inline_images(para, image_map: dict) -> str:
    """Extract inline images from paragraph XML and return as HTML img tags."""
    html = ""
    try:
        # Look for drawing elements with image references
        from lxml import etree
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
        }
        for blip in para._element.findall('.//a:blip', nsmap):
            embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if embed and embed in image_map:
                blob = image_map[embed]
                b64 = base64.b64encode(blob).decode('ascii')
                html += f'<img src="data:image/png;base64,{b64}">'
    except Exception:
        pass
    return html


def _table_to_html(table) -> str:
    """Convert a docx table to an HTML table."""
    rows_html = []
    for r_idx, row in enumerate(table.rows):
        cells = []
        tag = "th" if r_idx == 0 else "td"
        for cell in row.cells:
            text = cell.text.strip()
            text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            cells.append(f"<{tag}>{text}</{tag}>")
        rows_html.append(f"<tr>{''.join(cells)}</tr>")
    return f"<table>{''.join(rows_html)}</table>"
